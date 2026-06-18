#!/usr/bin/env python3
"""生成照片排列 PPT 模板 — 8 张照片分 2 页，大小不一，美观排列。"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── 颜色方案 ──────────────────────────────────────────
BG_COLOR      = RGBColor(0xFA, 0xF8, 0xF5)  # 暖白背景
ACCENT        = RGBColor(0x3B, 0x52, 0x6B)  # 深蓝灰强调
PHOTO_FILL    = RGBColor(0xE8, 0xE4, 0xDF)  # 照片占位底色
PHOTO_BORDER  = RGBColor(0xD0, 0xCC, 0xC7)  # 照片边框
TEXT_COLOR    = RGBColor(0x3B, 0x52, 0x6B)  # 文字色
SHADOW_COLOR  = RGBColor(0xDD, 0xD9, 0xD4)  # 阴影
STRIP_COLOR   = RGBColor(0x7E, 0x9C, 0xB5)  # 顶部装饰条
SUBTLE_COLOR  = RGBColor(0xB0, 0xAC, 0xA7)  # 次要文字
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── 辅助函数 ──────────────────────────────────────────
def add_photo_placeholder(slide, left, top, width, height, num, shadow=True):
    """添加一个照片占位框（含阴影 + 虚线边框 + 编号标签）。"""
    # 阴影层
    if shadow:
        shadow_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches(0.06), top + Inches(0.06),
            width, height
        )
        shadow_shape.fill.solid()
        shadow_shape.fill.fore_color.rgb = SHADOW_COLOR
        shadow_shape.line.fill.background()
        shadow_shape.rotation = 0
        # 调圆角
        try:
            shadow_shape.adjustments[0] = 0.04
        except Exception:
            pass

    # 主体框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PHOTO_FILL
    shape.line.color.rgb = PHOTO_BORDER
    shape.line.width = Pt(1.5)
    try:
        shape.adjustments[0] = 0.04
    except Exception:
        pass

    # 居中文字
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"📷\n照片 {num}"
    run.font.size = Pt(14)
    run.font.color.rgb = SUBTLE_COLOR
    run.font.name = "Microsoft YaHei"
    tf.paragraphs[0].space_before = Pt(0)

    # 垂直居中
    try:
        tf.word_wrap = True
        # 用 anchor 属性设置垂直居中
        from pptx.oxml.ns import qn
        txBody = shape._element.txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    except Exception:
        pass

    return shape


def add_slide_title(slide, text, sub_text=None):
    """添加页面标题 + 装饰条."""
    # 顶部装饰条
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.06)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = STRIP_COLOR
    strip.line.fill.background()

    # 标题
    left_margin = Inches(0.8)
    title_box = slide.shapes.add_textbox(
        left_margin, Inches(0.35),
        Inches(11), Inches(0.55)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Microsoft YaHei"

    if sub_text:
        sub_box = slide.shapes.add_textbox(
            left_margin, Inches(0.85),
            Inches(11), Inches(0.35)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = sub_text
        run2.font.size = Pt(13)
        run2.font.color.rgb = SUBTLE_COLOR
        run2.font.name = "Microsoft YaHei"


def set_slide_bg(slide, color):
    """设置幻灯片背景色."""
    from pptx.oxml.ns import qn
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── 第 1 页：非对称布局（5 张照片） ──────────────────
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
set_slide_bg(slide1, BG_COLOR)
add_slide_title(slide1, "精选记忆", "不对称布局 · 5 张照片")

MARGIN = Inches(0.8)
GAP = Inches(0.2)
TOP_OFFSET = Inches(1.45)

# 左侧大图 (照片 1) — 占左侧约 58%
big_w = Inches(6.8)
big_h = Inches(5.3)
big_left = MARGIN
big_top = TOP_OFFSET

add_photo_placeholder(slide1, big_left, big_top, big_w, big_h, 1)

# 右侧 4 张小图 — 2×2 网格
small_w = Inches(2.85)
small_h = Inches(2.55)
right_origin_x = Inches(8.05)  # 大图右边的起始位置
right_origin_y = TOP_OFFSET

positions = [
    (right_origin_x, right_origin_y),                        # 照片 2 — 右上
    (right_origin_x + small_w + GAP, right_origin_y),        # 照片 3 — 右上右
    (right_origin_x, right_origin_y + small_h + GAP),        # 照片 4 — 右下
    (right_origin_x + small_w + GAP, right_origin_y + small_h + GAP),  # 照片 5
]

for i, (lx, ly) in enumerate(positions):
    add_photo_placeholder(slide1, lx, ly, small_w, small_h, i + 2)


# ── 第 2 页：横幅布局（3 张照片） ──────────────────────
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_COLOR)
add_slide_title(slide2, "更多瞬间", "横幅与并排布局 · 3 张照片")

# 顶部宽幅横图 (照片 6)
wide_w = Inches(11.6)
wide_h = Inches(2.7)
wide_left = MARGIN
wide_top = TOP_OFFSET
add_photo_placeholder(slide2, wide_left, wide_top, wide_w, wide_h, 6)

# 下方两张竖图 (照片 7, 8)
bottom_top = wide_top + wide_h + GAP
bottom_w = Inches(5.7)
bottom_h = Inches(3.0)
bottom_left_1 = MARGIN
bottom_left_2 = MARGIN + bottom_w + GAP

add_photo_placeholder(slide2, bottom_left_1, bottom_top, bottom_w, bottom_h, 7)
add_photo_placeholder(slide2, bottom_left_2, bottom_top, bottom_w, bottom_h, 8)


# ── 第 3 页（可选）：单页全 8 张概览 ──────────────────
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_COLOR)
add_slide_title(slide3, "全览", "8 张照片 · 砖石布局")

# Masonry 风格 — 3 列高度不一
col_gap = Inches(0.18)
row_gap = Inches(0.18)
m_left = MARGIN
m_top = TOP_OFFSET

# 列宽度
col_w = [
    Inches(3.6),   # 左列
    Inches(4.3),   # 中列（略宽）
    Inches(3.6),   # 右列
]

# 布局定义: [(列, 高_inches), ...]
layout = [
    (0, 2.3),   # 照片 1 — 左列上方
    (1, 3.0),   # 照片 2 — 中列上方（较高）
    (2, 2.0),   # 照片 3 — 右列上方（较矮）
    (0, 2.8),   # 照片 4 — 左列下方
    (1, 2.0),   # 照片 5 — 中列下方
    (2, 3.2),   # 照片 6 — 右列下方
]

# 实际计算位置
col_x = [m_left, m_left + col_w[0] + col_gap, m_left + col_w[0] + col_gap + col_w[1] + col_gap]
col_y = [m_top, m_top, m_top]  # 当前每列的 y 偏移

for i, (col, h) in enumerate(layout):
    lx = col_x[col]
    ly = col_y[col]
    w = col_w[col]
    add_photo_placeholder(slide3, lx, ly, w, Inches(h), i + 1)
    col_y[col] = ly + Inches(h) + row_gap

# 剩余 2 张放在底行
bottom_row_y = max(col_y) + Inches(0.1)
bottom_h_final = Inches(2.0)
# 照片 7 — 横跨左+中部分
add_photo_placeholder(
    slide3,
    m_left,
    bottom_row_y,
    col_w[0] + col_gap + col_w[1],
    bottom_h_final,
    7
)
# 照片 8 — 右侧
add_photo_placeholder(
    slide3,
    col_x[2],
    bottom_row_y,
    col_w[2],
    bottom_h_final,
    8
)


# ── 保存 ──────────────────────────────────────────────
output_path = "/Users/harvey/Desktop/照片排列模板.pptx"
prs.save(output_path)
print(f"✅ PPT 已保存至: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
